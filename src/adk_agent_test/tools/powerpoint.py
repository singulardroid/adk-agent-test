"""PowerPoint tool: create a local .pptx and add content (e.g. final output)."""

import os
import sys
from pathlib import Path
from typing import Annotated

from ..config import is_mock_tools
from ..logging_config import get_logger
from ._mockable import mockable

logger = get_logger(__name__)

RESULTS_DIR = "results"


def _resolve_output_path(output_path: str) -> Path:
    """Resolve output path: relative paths are saved under ./results/."""
    path = Path(output_path)
    if not path.is_absolute():
        path = Path(RESULTS_DIR) / path.name
    return path


def mock_create_powerpoint(
    content: Annotated[str, "Text to put in the presentation (e.g. final output or summary)"],
    output_path: Annotated[str, "Path where to save the .pptx file"] = "output.pptx",
    title: Annotated[str, "Title for the first slide"] = "Summary",
) -> str:
    """Mock create_powerpoint: no file I/O."""
    resolved = _resolve_output_path(output_path)
    logger.info(
        "[MOCK] mock_create_powerpoint called output_path=%s title=%r content_len=%d",
        str(resolved),
        (title or "")[:40],
        len(content or ""),
    )
    print("[MOCK] mock_create_powerpoint called", flush=True)
    return f"Saved mock presentation to {resolved}"


@mockable
def create_powerpoint(
    content: Annotated[str, "Text to put in the presentation (e.g. final output or summary)"],
    output_path: Annotated[str, "Path where to save the .pptx file"] = "output.pptx",
    title: Annotated[str, "Title for the first slide"] = "Summary",
) -> str:
    """Create a local PowerPoint file and add the given content to it."""
    if is_mock_tools():
        msg = "ADK_MOCK_TOOLS=1 but real create_powerpoint() was invoked; mock path should have been used."
        logger.error("%s", msg)
        print(msg, file=sys.stderr, flush=True)
        os._exit(1)
    logger.info("[REAL] create_powerpoint() called output_path=%s", output_path)
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as e:
        logger.exception("python-pptx not available: %s", e)
        return f"Error: python-pptx is not installed. Install with: pip install python-pptx\n{e!s}"
    path = _resolve_output_path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title or "Summary"
    # Content slide
    content_slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(content_slide_layout)
    left = Inches(0.5)
    top = Inches(0.5)
    width = prs.slide_width - Inches(1)
    height = prs.slide_height - Inches(1)
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (content or "").strip() or "(No content)"
    p.font.size = Pt(12)
    try:
        prs.save(str(path))
    except (OSError, IOError) as e:
        logger.exception("Failed to save presentation to %s: %s", path, e)
        return f"Error saving presentation to {path}: {e!s}"
    logger.info("create_powerpoint() saved to %s", path)
    return f"Saved presentation to {path.absolute()}"
