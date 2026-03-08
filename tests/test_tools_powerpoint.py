"""Tests for PowerPoint tool with mock and real implementation."""

import importlib
import os
from pathlib import Path

import pytest

import adk_agent_test.tools.powerpoint as ppt_module
from adk_agent_test.tools.powerpoint import create_powerpoint


# ----- ADK_MOCK_TOOLS=1: decorator uses mock_* implementation -----


def test_adk_mock_tools_create_powerpoint_returns_mock_result() -> None:
    """When ADK_MOCK_TOOLS=1, create_powerpoint() uses mock_create_powerpoint and does not write a file."""
    os.environ["ADK_MOCK_TOOLS"] = "1"
    try:
        importlib.reload(ppt_module)
        result = ppt_module.create_powerpoint(
            content="My summary here",
            output_path="any.pptx",
            title="Report",
        )
        assert "[MOCK]" in result or "mock" in result.lower()
        assert "any.pptx" in result
        assert "Saved mock presentation" in result
        assert not Path("any.pptx").exists()
    finally:
        os.environ.pop("ADK_MOCK_TOOLS", None)
        importlib.reload(ppt_module)


# ----- Real implementation: creates .pptx on disk -----


def test_create_powerpoint_real_creates_file_with_content(tmp_path: Path) -> None:
    """When ADK_MOCK_TOOLS=0, create_powerpoint() writes a .pptx and content is in the slides."""
    os.environ["ADK_MOCK_TOOLS"] = "0"
    try:
        importlib.reload(ppt_module)
        out = tmp_path / "test_out.pptx"
        content = "Hello from the agent. This is the final output."
        result = ppt_module.create_powerpoint(
            content=content,
            output_path=str(out),
            title="Summary",
        )
        assert "Saved presentation to" in result or "test_out.pptx" in result
        assert out.exists()
        from pptx import Presentation

        prs = Presentation(str(out))
        assert len(prs.slides) >= 2
        # First slide: title
        assert prs.slides[0].shapes.title.text == "Summary"
        # Second slide: content
        text_frame = prs.slides[1].shapes[0].text_frame
        full_text = "".join(p.text for p in text_frame.paragraphs)
        assert "Hello from the agent" in full_text
        assert "final output" in full_text
    finally:
        os.environ.pop("ADK_MOCK_TOOLS", None)
        importlib.reload(ppt_module)
